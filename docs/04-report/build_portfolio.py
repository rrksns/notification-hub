# Notification Hub 포트폴리오 PPTX를 Claude Code 문서 톤(아이보리+코랄)으로 렌더링하는 전용 빌드 스크립트
"""실제 도형 다이어그램과 네이티브 차트로 슬라이드를 그린다.

Usage:
    /tmp/pptx-venv/bin/python build_portfolio.py [output.pptx]
"""

import sys
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

# ── 디자인 토큰 (Claude Code 문서 사이트 색감) ─────────────────────────────────
IVORY      = RGBColor(0xFA, 0xF9, 0xF5)   # 배경
CARD       = RGBColor(0xF2, 0xEF, 0xE7)   # 카드/패널
CARD_DK    = RGBColor(0xEA, 0xE6, 0xDA)   # 카드(진한)
BORDER     = RGBColor(0xE0, 0xDA, 0xCC)   # 보더
CORAL      = RGBColor(0xD9, 0x77, 0x57)   # 시그니처 코랄
CORAL_DK   = RGBColor(0xC2, 0x61, 0x3F)   # 코랄 강조
CORAL_SOFT = RGBColor(0xF4, 0xE2, 0xD7)   # 코랄 틴트
INK        = RGBColor(0x20, 0x1F, 0x1B)   # 다크 섹션 배경
INK_SOFT   = RGBColor(0x2C, 0x2A, 0x24)   # 다크 카드
TEXT       = RGBColor(0x33, 0x31, 0x2B)   # 본문 텍스트
MUTED      = RGBColor(0x6E, 0x6A, 0x60)   # 보조 텍스트
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
SLATE      = RGBColor(0x5E, 0x6B, 0x66)   # 보조 톤(다이어그램)
SLATE_SOFT = RGBColor(0xE3, 0xE7, 0xE3)

FONT = "Apple SD Gothic Neo"

EMU_IN = 914400


def IN(v):
    return Emu(int(v * EMU_IN))


# ── 기본 헬퍼 ─────────────────────────────────────────────────────────────────
def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def bg(slide, color):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color


def _no_shadow(shape):
    shape.shadow.inherit = False


def rect(slide, l, t, w, h, fill, line=None, line_w=1.0):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, IN(l), IN(t), IN(w), IN(h))
    _no_shadow(sp)
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    return sp


def rrect(slide, l, t, w, h, fill, line=None, line_w=1.0, radius=0.12):
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, IN(l), IN(t), IN(w), IN(h))
    _no_shadow(sp)
    try:
        sp.adjustments[0] = radius
    except Exception:
        pass
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    return sp


def oval(slide, l, t, w, h, fill, line=None, line_w=1.0):
    sp = slide.shapes.add_shape(MSO_SHAPE.OVAL, IN(l), IN(t), IN(w), IN(h))
    _no_shadow(sp)
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(line_w)
    return sp


def can(slide, l, t, w, h, fill, line=None):
    sp = slide.shapes.add_shape(MSO_SHAPE.CAN, IN(l), IN(t), IN(w), IN(h))
    _no_shadow(sp)
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
    return sp


def arrow_r(slide, l, t, w, h, fill=CORAL):
    sp = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, IN(l), IN(t), IN(w), IN(h))
    _no_shadow(sp)
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    sp.line.fill.background()
    try:
        sp.adjustments[0] = 0.55
        sp.adjustments[1] = 0.55
    except Exception:
        pass
    return sp


def arrow_d(slide, l, t, w, h, fill=CORAL):
    sp = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, IN(l), IN(t), IN(w), IN(h))
    _no_shadow(sp)
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    sp.line.fill.background()
    return sp


def line(slide, x1, y1, x2, y2, color=CORAL, w=1.6, arrow=True, dash=False):
    cn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, IN(x1), IN(y1), IN(x2), IN(y2))
    cn.line.color.rgb = color
    cn.line.width = Pt(w)
    ln = cn.line._get_or_add_ln()
    if arrow:
        tail = ln.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'})
        ln.append(tail)
    if dash:
        d = ln.makeelement(qn('a:prstDash'), {'val': 'dash'})
        ln.insert(0, d)
    return cn


def text(slide, s, l, t, w, h, size=14, bold=False, color=TEXT,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT, spacing=None):
    tb = slide.shapes.add_textbox(IN(l), IN(t), IN(w), IN(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Pt(2); tf.margin_right = Pt(2)
    tf.margin_top = Pt(1); tf.margin_bottom = Pt(1)
    p = tf.paragraphs[0]
    p.alignment = align
    if spacing:
        p.line_spacing = spacing
    r = p.add_run(); r.text = s
    r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = color; r.font.name = font
    return tb


def node(slide, l, t, w, h, title, fill=WHITE, line=BORDER, tcolor=TEXT,
         size=12, bold=True, sub=None, sub_size=9, sub_color=MUTED, radius=0.14):
    rrect(slide, l, t, w, h, fill, line=line, line_w=1.25, radius=radius)
    if sub:
        text(slide, title, l, t + h * 0.12, w, h * 0.5, size=size, bold=bold,
             color=tcolor, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.BOTTOM)
        text(slide, sub, l, t + h * 0.5, w, h * 0.42, size=sub_size, bold=False,
             color=sub_color, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)
    else:
        text(slide, title, l, t, w, h, size=size, bold=bold, color=tcolor,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def chip(slide, l, t, w, h, s, fill=CORAL_SOFT, tcolor=CORAL_DK, size=9, bold=True):
    rrect(slide, l, t, w, h, fill, line=None, radius=0.5)
    text(slide, s, l, t, w, h, size=size, bold=bold, color=tcolor,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def bullets(slide, items, l, t, w, h, size=13, gap=10, marker="▸"):
    tb = slide.shapes.add_textbox(IN(l), IN(t), IN(w), IN(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(0); tf.margin_right = Pt(0)
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(0 if i == 0 else gap)
        p.space_after = Pt(0)
        p.line_spacing = 1.12
        rm = p.add_run(); rm.text = marker + "  "
        rm.font.size = Pt(size); rm.font.bold = True
        rm.font.color.rgb = CORAL; rm.font.name = FONT
        rt = p.add_run(); rt.text = it
        rt.font.size = Pt(size); rt.font.color.rgb = TEXT; rt.font.name = FONT
    return tb


W, H = 13.333, 7.5
ACCENT_W = 0.16


def content_header(slide, kicker, title, key_msg):
    """좌측 코랄 바 + 키커 + 타이틀 + 핵심 메시지 밴드. content 시작 y를 반환."""
    bg(slide, IVORY)
    rect(slide, 0, 0, ACCENT_W, H, CORAL)
    text(slide, kicker, 0.62, 0.46, 11.0, 0.34, size=11.5, bold=True,
         color=CORAL_DK, align=PP_ALIGN.LEFT)
    text(slide, title, 0.6, 0.78, 12.2, 0.62, size=27, bold=True,
         color=TEXT, align=PP_ALIGN.LEFT)
    # 핵심 메시지 밴드
    by, bh = 1.62, 0.92
    rrect(slide, 0.62, by, 12.1, bh, CARD, line=BORDER, line_w=1.0, radius=0.10)
    rect(slide, 0.62, by, 0.10, bh, CORAL)
    tb = slide.shapes.add_textbox(IN(0.92), IN(by), IN(11.6), IN(bh))
    tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.line_spacing = 1.1
    r1 = p.add_run(); r1.text = "핵심   "
    r1.font.size = Pt(12); r1.font.bold = True; r1.font.color.rgb = CORAL_DK; r1.font.name = FONT
    r2 = p.add_run(); r2.text = key_msg
    r2.font.size = Pt(12.5); r2.font.bold = True; r2.font.color.rgb = TEXT; r2.font.name = FONT
    return 2.82


def page_no(slide, n):
    text(slide, str(n), W - 0.7, H - 0.45, 0.4, 0.3, size=10, color=MUTED,
         align=PP_ALIGN.RIGHT)


# ── 슬라이드 빌더 ─────────────────────────────────────────────────────────────
def s_title(prs):
    s = blank(prs)
    bg(s, IVORY)
    rect(s, 0, 0, ACCENT_W, H, CORAL)
    # 로고 마크
    rrect(s, 0.9, 0.95, 0.66, 0.66, CORAL, radius=0.28)
    text(s, "N", 0.9, 0.95, 0.66, 0.66, size=26, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, "B2B SaaS · EVENT-DRIVEN MSA", 0.92, 2.55, 9.0, 0.4, size=13,
         bold=True, color=CORAL_DK)
    text(s, "Notification Hub", 0.86, 2.95, 11.0, 1.2, size=58, bold=True, color=TEXT)
    text(s, "멀티테넌트 알림 발송 플랫폼 · Clean Architecture 기반 Event-Driven MSA",
         0.92, 4.32, 10.5, 0.6, size=17, color=MUTED)
    rect(s, 0.92, 5.05, 2.6, 0.045, CORAL)
    text(s, "이설      |      2026-06-05", 0.92, 6.55, 8.0, 0.4, size=13, color=TEXT)
    # 우측 장식: 살짝 겹친 카드 3장
    for i, c in enumerate([CARD_DK, CARD, CORAL_SOFT]):
        rrect(s, 9.7 + i * 0.28, 1.5 + i * 0.55, 2.7, 1.5, c, line=BORDER, line_w=1.0, radius=0.10)
    text(s, "Email", 10.05, 1.72, 2.0, 0.4, size=12, bold=True, color=MUTED)
    text(s, "SMS", 10.33, 2.27, 2.0, 0.4, size=12, bold=True, color=MUTED)
    text(s, "Push", 10.61, 2.82, 2.0, 0.4, size=12, bold=True, color=CORAL_DK)


def s_agenda(prs):
    s = blank(prs)
    bg(s, IVORY)
    rect(s, 0, 0, ACCENT_W, H, CORAL)
    text(s, "AGENDA", 0.62, 0.7, 6.0, 0.4, size=13, bold=True, color=CORAL_DK)
    text(s, "목차", 0.6, 1.05, 6.0, 0.8, size=34, bold=True, color=TEXT)
    items = [
        "프로젝트 개요 & 문제 정의",
        "시스템 아키텍처",
        "핵심 플로우 — 이벤트 드리븐 파이프라인",
        "신뢰성 엔지니어링",
        "운영 & 품질",
        "성과 요약 & 배운 점",
    ]
    top, step = 2.25, 0.78
    for i, it in enumerate(items):
        y = top + i * step
        rrect(s, 0.92, y, 0.56, 0.56, CORAL, radius=0.22)
        text(s, str(i + 1), 0.92, y, 0.56, 0.56, size=18, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        text(s, it, 1.72, y, 10.0, 0.56, size=18, bold=False, color=TEXT,
             anchor=MSO_ANCHOR.MIDDLE)


def s_section(prs, num, title):
    s = blank(prs)
    bg(s, INK)
    rect(s, 0, 0, ACCENT_W, H, CORAL)
    text(s, num, 0.85, 2.0, 4.0, 2.2, size=120, bold=True, color=CORAL,
         anchor=MSO_ANCHOR.MIDDLE)
    text(s, title, 3.6, 2.9, 9.0, 1.2, size=40, bold=True, color=WHITE,
         anchor=MSO_ANCHOR.MIDDLE)
    rect(s, 3.65, 4.15, 3.0, 0.05, CORAL)


# S4 ── 문제 정의: 각자 구축 vs Hub 위임
def s_problem(prs):
    s = blank(prs)
    y0 = content_header(
        s, "01 — 프로젝트 개요",
        "프로젝트 개요 & 문제 정의",
        "고객사가 알림 인프라를 직접 구축할 필요 없이, API 한 번으로 Email/SMS/Push 발송을 위임하는 B2B SaaS 플랫폼")
    bullets(s, [
        "문제: 모든 서비스가 알림 발송·재시도·통계 인프라를 중복 구축 → 비용·운영 부담",
        "해결: 멀티테넌트 허브가 발송·재시도·집계를 대행, 고객사는 REST API만 호출",
        "테넌트 격리: JWT 클레임 기반 X-Tenant-Id 주입으로 고객사 간 데이터 완전 분리",
        "규모: Java 17 · Spring Boot 3.2 · 6개 MSA · Kafka 이벤트 파이프라인",
    ], 0.7, y0 + 0.15, 5.0, 3.9, size=12.5, gap=12)

    # 다이어그램 캔버스 (우측)
    cx, cy, cw = 6.05, y0 + 0.05, 6.6
    # 좌 패널: 기존(중복)
    text(s, "기존 — 서비스마다 각자 구축", cx, cy, cw / 2 - 0.1, 0.34, size=11.5,
         bold=True, color=MUTED, align=PP_ALIGN.CENTER)
    px = cx
    for i in range(3):
        yy = cy + 0.5 + i * 0.92
        node(s, px + 0.15, yy, 2.6, 0.74, "서비스 " + "ABC"[i], fill=WHITE, line=BORDER,
             sub="+ 알림·재시도·통계 직접 구현", size=12, sub_size=8.5, sub_color=CORAL_DK)
    # 구분선
    rect(s, cx + cw / 2 - 0.02, cy + 0.45, 0.02, 2.95, BORDER)

    # 우 패널: Hub 위임
    rx = cx + cw / 2 + 0.15
    text(s, "Notification Hub — 단일 허브 위임", rx - 0.15, cy, cw / 2, 0.34,
         size=11.5, bold=True, color=CORAL_DK, align=PP_ALIGN.CENTER)
    for i in range(3):
        yy = cy + 0.55 + i * 0.62
        node(s, rx, yy, 1.15, 0.46, "서비스 " + "ABC"[i], fill=CARD, line=BORDER, size=9.5)
        line(s, rx + 1.18, yy + 0.23, rx + 1.78, cy + 1.62, color=CORAL, w=1.4)
    # Hub
    rrect(s, rx + 1.8, cy + 1.18, 1.5, 0.9, CORAL, radius=0.16)
    text(s, "Notification\nHub", rx + 1.8, cy + 1.18, 1.5, 0.9, size=11.5, bold=True,
         color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # 채널
    chs = ["Email", "SMS", "Push"]
    for i, ch in enumerate(chs):
        line(s, rx + 2.55, cy + 2.1, rx + 1.95 + i * 0.62, cy + 2.55, color=CORAL, w=1.2, arrow=False)
        chip(s, rx + 1.62 + i * 0.62, cy + 2.55, 0.58, 0.34, ch, size=8.5)
    page_no(s, 1)


# S6 ── 아키텍처: Clean Arch 동심원 + MSA 토폴로지
def s_arch(prs):
    s = blank(prs)
    y0 = content_header(
        s, "02 — 시스템 아키텍처",
        "시스템 아키텍처",
        "6개 MSA + Clean Architecture(Port & Adapter)로 도메인 로직을 Spring/JPA 같은 인프라에서 완전히 분리")
    bullets(s, [
        "6개 서비스: Discovery · Gateway · User · Notification · Delivery · Analytics",
        "의존성 방향: presentation → application → domain (domain은 외부 의존 0)",
        "인프라 교체 시 도메인 변경 없이 어댑터만 교체 → ArchUnit으로 규칙 자동 검증",
        "폴리글랏 저장소: 트랜잭션 MySQL · 통계 MongoDB · 멱등성/카운터 Redis",
    ], 0.7, y0 + 0.15, 5.0, 3.9, size=12.5, gap=12)

    cx, cy = 6.05, y0 + 0.1
    # 동심원 Clean Architecture
    text(s, "Clean Architecture", cx, cy, 3.0, 0.32, size=11, bold=True,
         color=MUTED, align=PP_ALIGN.CENTER)
    rings = [
        (3.0, 3.0, CARD_DK, "presentation", TEXT, 9),
        (2.32, 2.32, SLATE_SOFT, "application", TEXT, 9),
        (1.64, 1.64, CORAL_SOFT, "infrastructure", CORAL_DK, 9),
        (0.96, 0.96, CORAL, "domain", WHITE, 11),
    ]
    ccx, ccy = cx + 1.5, cy + 0.5
    for w_, h_, col, lab, tc, fs in rings:
        oval(s, ccx - w_ / 2, ccy + 1.5 - h_ / 2, w_, h_, col, line=WHITE, line_w=1.5)
    for w_, h_, col, lab, tc, fs in rings:
        yy = ccy + 1.5 - h_ / 2
        text(s, lab, ccx - w_ / 2, yy + 0.04, w_, 0.3, size=fs, bold=True,
             color=tc, align=PP_ALIGN.CENTER) if lab != "domain" else \
            text(s, lab, ccx - w_ / 2, ccy + 1.5 - 0.15, w_, 0.3, size=fs, bold=True,
                 color=tc, align=PP_ALIGN.CENTER)

    # MSA 토폴로지 (우측)
    mx = cx + 3.25
    text(s, "MSA 토폴로지", mx, cy, 3.2, 0.32, size=11, bold=True,
         color=MUTED, align=PP_ALIGN.CENTER)
    node(s, mx, cy + 0.4, 3.2, 0.5, "API Gateway  (JWT 인증·라우팅)", fill=CORAL,
         line=None, tcolor=WHITE, size=10.5)
    svcs = ["User", "Notification", "Delivery", "Analytics", "Discovery", "(Eureka)"]
    for i, sv in enumerate(svcs):
        r, c = divmod(i, 2)
        node(s, mx + c * 1.65, cy + 1.1 + r * 0.66, 1.5, 0.52,
             sv, fill=WHITE if sv != "(Eureka)" else CARD, line=BORDER, size=10)
    # gateway → services 연결
    line(s, mx + 1.6, cy + 0.9, mx + 1.6, cy + 1.08, color=CORAL, w=1.4, arrow=True)
    page_no(s, 2)


# S8 ── 이벤트 드리븐 파이프라인 (가로 풀폭)
def s_pipeline(prs):
    s = blank(prs)
    y0 = content_header(
        s, "03 — 핵심 플로우",
        "핵심 플로우 — 이벤트 드리븐",
        "Kafka 기반 비동기 파이프라인으로 '접수 → 발송 → 집계'를 분리해, 발송 지연이 알림 접수 응답을 막지 않도록 설계")
    # 파이프라인 밴드
    py = y0 + 0.25
    bh = 0.78
    # 노드 정의: (x, w, label, kind)  kind: box | gw | kafka | sink
    nodes = [
        (0.7, 1.25, "Client", "box"),
        (2.25, 1.35, "API\nGateway", "gw"),
        (3.95, 1.55, "Notification\n접수", "box"),
        (5.85, 1.4, "notifications", "kafka"),
        (7.65, 1.5, "Delivery\n발송", "box"),
        (9.5, 1.4, "delivery-results", "kafka"),
        (11.3, 1.45, "Analytics\n집계", "box"),
    ]
    centers = []
    for x, w_, lab, kind in nodes:
        cyc = py + bh / 2
        if kind == "kafka":
            can(s, x, py - 0.05, w_, bh + 0.1, CORAL)
            text(s, "Kafka", x, py + 0.0, w_, 0.3, size=8, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)
            text(s, lab, x, py + 0.26, w_, 0.5, size=9.5, bold=True, color=WHITE,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        elif kind == "gw":
            node(s, x, py, w_, bh, lab, fill=INK_SOFT, line=None, tcolor=WHITE, size=10.5)
        else:
            node(s, x, py, w_, bh, lab, fill=WHITE, line=BORDER, size=10.5)
        centers.append((x, x + w_, cyc))
    # 화살표
    for i in range(len(centers) - 1):
        x_end = centers[i][1]
        x_next = centers[i + 1][0]
        arrow_r(s, x_end + 0.04, py + bh / 2 - 0.13, (x_next - x_end) - 0.08, 0.26)
    # 응답 라벨
    text(s, "↩ 즉시 PUBLISHED 응답", 3.95, py + bh + 0.12, 2.4, 0.3, size=9.5,
         bold=True, color=CORAL_DK)
    # 저장소 태그
    chip(s, 3.95, py - 0.42, 1.0, 0.3, "MySQL", fill=SLATE_SOFT, tcolor=SLATE)
    chip(s, 11.3, py - 0.42, 1.45, 0.3, "MongoDB · Redis", fill=SLATE_SOFT, tcolor=SLATE)

    # 하단 불릿 (2열)
    by = py + bh + 0.62
    left = [
        "접수: 멱등성 체크 → MySQL 저장 → notifications 발행 → 즉시 응답",
        "발송: Kafka 소비 → 채널별 발송 → delivery-results 발행",
    ]
    right = [
        "집계: 결과 소비 → MongoDB 일별 통계 + Redis 실시간 카운터",
        "각 단계는 토픽으로만 통신 → 독립 배포·확장 가능",
    ]
    bullets(s, left, 0.7, by, 5.9, 1.3, size=11.5, gap=9)
    bullets(s, right, 6.85, by, 5.9, 1.3, size=11.5, gap=9)
    page_no(s, 3)


# S9 ── 신뢰성: 재시도/DLQ + Circuit Breaker
def s_reliability(prs):
    s = blank(prs)
    y0 = content_header(
        s, "04 — 신뢰성 엔지니어링",
        "신뢰성 엔지니어링",
        "멱등성·Circuit Breaker·재시도/DLQ·원자적 집계로 '메시지 유실 없고 중복 없는' 발송을 보장")
    bullets(s, [
        "멱등성 이중 방어: 접수단 Redis 키 + 발송단 notificationId 중복 체크",
        "Kafka 발행 신뢰성: fire-and-forget 대신 동기 확인(.get 5초) + 실패 시 예외 전파",
        "장애 격리: Resilience4j Circuit Breaker(실패율 50% 초과 시 차단)",
        "최종 실패 보존: 3회 재시도 후 DLQ 격리 + $inc·upsert 원자적 집계",
    ], 0.7, y0 + 0.15, 5.0, 3.9, size=12.5, gap=12)

    cx, cy = 6.05, y0 + 0.2
    # 재시도/DLQ 흐름
    text(s, "재시도 → DLQ 격리", cx, cy, 6.5, 0.32, size=11, bold=True,
         color=MUTED, align=PP_ALIGN.LEFT)
    flow = [("발송 실패", CARD), ("retry 1s", WHITE), ("retry 2s", WHITE),
            ("retry 4s", WHITE), ("DLQ", CORAL)]
    fx, fy, fw, fh = cx, cy + 0.42, 1.12, 0.62
    for i, (lab, col) in enumerate(flow):
        x = fx + i * (fw + 0.18)
        tc = WHITE if col == CORAL else TEXT
        node(s, x, fy, fw, fh, lab, fill=col, line=None if col != WHITE else BORDER,
             tcolor=tc, size=10)
        if i < len(flow) - 1:
            arrow_r(s, x + fw + 0.01, fy + fh / 2 - 0.1, 0.16, 0.2)

    # Circuit Breaker 상태 머신
    sy = cy + 1.55
    text(s, "Circuit Breaker 상태 전이", cx, sy, 6.5, 0.32, size=11, bold=True,
         color=MUTED, align=PP_ALIGN.LEFT)
    states = [("CLOSED", SLATE), ("OPEN", CORAL), ("HALF_OPEN", CORAL_DK)]
    d = 1.35
    sxs = [cx + 0.1, cx + 2.55, cx + 5.0]
    syy = sy + 0.6
    for (lab, col), sx in zip(states, sxs):
        oval(s, sx, syy, d, d, col)
        text(s, lab, sx, syy, d, d, size=10.5, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    line(s, sxs[0] + d, syy + d / 2, sxs[1], syy + d / 2, color=TEXT, w=1.4)
    text(s, "실패율 50%↑", sxs[0] + d + 0.02, syy + d / 2 - 0.34, 1.1, 0.3, size=8.5, color=MUTED)
    line(s, sxs[1] + d, syy + d / 2, sxs[2], syy + d / 2, color=TEXT, w=1.4)
    text(s, "타임아웃 후", sxs[1] + d + 0.05, syy + d / 2 - 0.34, 1.1, 0.3, size=8.5, color=MUTED)
    line(s, sxs[2] + d / 2, syy + d, sxs[0] + d / 2, syy + d + 0.0, color=SLATE, w=1.2, arrow=True)
    text(s, "성공 시 복구 →", sxs[1] - 0.2, syy + d + 0.05, 1.6, 0.3, size=8.5, color=SLATE)
    page_no(s, 4)


# S11 ── 운영 & 품질: 불릿 + 네이티브 커버리지 차트
def s_ops(prs):
    s = blank(prs)
    y0 = content_header(
        s, "05 — 운영 & 품질",
        "운영 & 품질",
        "모니터링·컨테이너 오케스트레이션·CI/CD·높은 테스트 커버리지로 '운영 가능한' 수준까지 완성")
    bullets(s, [
        "관측성: Prometheus + Grafana(발송량·5xx·Lag) + Zipkin 분산 트레이싱",
        "배포: Docker Compose(로컬) · Kubernetes(HPA 2~10 Pod) · Terraform IaC",
        "CI/CD: GitHub Actions 전체 테스트 + 7개 서비스 Docker 이미지 매트릭스 빌드",
        "테스트: 도메인+애플리케이션 계층 커버리지 83.5% ~ 94.7%",
    ], 0.7, y0 + 0.15, 5.4, 3.9, size=12.5, gap=12)

    # 네이티브 막대 차트
    cx, cy, cw, ch = 6.4, y0 + 0.05, 6.3, 3.7
    text(s, "계층별 테스트 커버리지 (%)", cx, cy, cw, 0.34, size=12, bold=True,
         color=TEXT, align=PP_ALIGN.CENTER)
    cd = CategoryChartData()
    cd.categories = ["user", "notification", "delivery", "analytics"]
    cd.add_series("커버리지", (90.0, 83.5, 93.4, 94.7))
    gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED,
                            IN(cx), IN(cy + 0.4), IN(cw), IN(ch - 0.4), cd)
    chart = gf.chart
    chart.has_legend = False
    chart.has_title = False
    plot = chart.plots[0]
    plot.gap_width = 70
    plot.has_data_labels = True
    dl = plot.data_labels
    dl.number_format = '0.0'
    dl.number_format_is_linked = False
    dl.font.size = Pt(11); dl.font.bold = True; dl.font.color.rgb = INK
    dl.font.name = FONT
    series = plot.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = CORAL
    va = chart.value_axis
    va.minimum_scale = 0; va.maximum_scale = 100
    va.has_major_gridlines = True
    va.major_gridlines.format.line.color.rgb = BORDER
    va.major_gridlines.format.line.width = Pt(0.5)
    va.tick_labels.font.size = Pt(9); va.tick_labels.font.color.rgb = MUTED
    va.format.line.color.rgb = BORDER
    ca = chart.category_axis
    ca.tick_labels.font.size = Pt(10); ca.tick_labels.font.color.rgb = TEXT
    ca.tick_labels.font.name = FONT
    ca.format.line.color.rgb = BORDER
    page_no(s, 5)


# S12 ── 성과 요약 & 배운 점 (다크)
def s_conclusion(prs):
    s = blank(prs)
    bg(s, INK)
    rect(s, 0, 0, ACCENT_W, H, CORAL)
    text(s, "SO WHAT", 0.7, 0.7, 6.0, 0.4, size=13, bold=True, color=CORAL)
    text(s, "성과 요약 & 배운 점", 0.68, 1.05, 11.0, 0.8, size=32, bold=True, color=WHITE)
    items = [
        ("성과", "Clean Architecture + Event-Driven MSA로 확장·교체·테스트가 용이한 알림 플랫폼 설계·구현"),
        ("기술 학습", "분산 환경의 멱등성/재시도/DLQ, Circuit Breaker, 원자적 집계로 신뢰성 확보"),
        ("운영 학습", "관측성(Prometheus/Zipkin)과 오케스트레이션(K8s/HPA)으로 운영 관점까지 경험"),
        ("개선 방향", "실제 채널(SendGrid/Twilio/FCM) 연동, Outbox 패턴 도입, E2E 자동화 확대"),
    ]
    top, step = 2.3, 0.92
    for i, (k, v) in enumerate(items):
        y = top + i * step
        oval(s, 0.72, y + 0.06, 0.34, 0.34, CORAL)
        text(s, "✓", 0.72, y + 0.04, 0.34, 0.34, size=13, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        text(s, k, 1.3, y, 2.0, 0.45, size=15, bold=True, color=CORAL, anchor=MSO_ANCHOR.MIDDLE)
        text(s, v, 3.0, y - 0.03, 9.6, 0.55, size=13.5, color=RGBColor(0xE6, 0xE2, 0xD8),
             anchor=MSO_ANCHOR.MIDDLE)
    rrect(s, 0.7, H - 1.15, 11.95, 0.62, INK_SOFT, line=CORAL, line_w=1.0, radius=0.2)
    text(s, "GitHub 저장소와 설계 문서(README · kafka-redis.md)로 코드와 의사결정 과정을 확인하실 수 있습니다.",
         0.9, H - 1.15, 11.6, 0.62, size=12.5, bold=True, color=RGBColor(0xEC, 0xE8, 0xDE),
         anchor=MSO_ANCHOR.MIDDLE)


def s_qa(prs):
    s = blank(prs)
    bg(s, INK)
    rect(s, 0, H - 0.45, W, 0.45, CORAL)
    rrect(s, W / 2 - 0.9, 1.95, 1.8, 1.8, CORAL, radius=0.3)
    text(s, "?", W / 2 - 0.9, 1.9, 1.8, 1.8, size=80, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, "Q & A", 0, 3.95, W, 1.0, size=46, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER)
    text(s, "감사합니다", 0, 5.0, W, 0.6, size=20, color=RGBColor(0xCB, 0xC6, 0xBA),
         align=PP_ALIGN.CENTER)
    text(s, "이설 · rrksns@nate.com", 0, 5.65, W, 0.5, size=14, color=MUTED,
         align=PP_ALIGN.CENTER)


# 발표자 노트
NOTES = {
    "problem": "도입부에서 '왜 만들었는가'를 공감시키는 슬라이드입니다. 알림은 거의 모든 서비스가 필요로 하지만 재시도·통계·멀티채널을 제대로 구현하려면 부담이 큽니다. 채널 실제 연동(SendGrid/Twilio/FCM)은 스텁이며 아키텍처 설계에 집중했다는 점을 솔직히 밝힙니다.",
    "arch": "핵심은 '왜 Clean Architecture인가'입니다. 도메인이 프레임워크를 모르기 때문에 인프라 없이 단위 테스트가 쉽고 기술 교체에 강합니다. ArchUnit으로 'domain이 infrastructure를 import하면 빌드 실패'하도록 강제해 설계가 시간이 지나도 무너지지 않게 했습니다.",
    "pipeline": "동기 호출이 아니라 이벤트로 단계를 분리한 이유를 설명합니다. 발송 채널이 느리거나 장애여도 고객사는 접수 응답을 즉시 받습니다. 파티션 3개로 발송을 병렬 처리하고, Zipkin으로 Kafka 구간까지 추적합니다.",
    "reliability": "기술 면접관이 가장 주목할 부분입니다. 멱등성으로 중복을 흡수하고, 재시도·DLQ로 유실을 막고, Circuit Breaker로 장애 전파를 차단했습니다. 통계 집계 동시성 버그를 read-modify-write에서 MongoDB 원자적 연산으로 바꾼 경험도 함께 언급합니다.",
    "ops": "기능 구현으로 끝내지 않고 '운영'까지 고려했음을 보여줍니다. 커버리지 수치는 도메인·애플리케이션 계층 기준이며, 비즈니스 로직이 프레임워크와 분리되어 높은 커버리지가 가능했다는 점을 Clean Architecture 효과와 연결합니다.",
}


def main():
    out = sys.argv[1] if len(sys.argv) > 1 else "Notification_Hub_포트폴리오.pptx"
    prs = Presentation()
    prs.slide_width = IN(W)
    prs.slide_height = IN(H)

    s_title(prs)
    s_agenda(prs)
    s_section(prs, "01", "프로젝트 개요")
    s_problem(prs)
    s_section(prs, "02", "시스템 아키텍처")
    s_arch(prs)
    s_section(prs, "03", "핵심 플로우 & 신뢰성")
    s_pipeline(prs)
    s_reliability(prs)
    s_section(prs, "04", "운영 & 성과")
    s_ops(prs)
    s_conclusion(prs)
    s_qa(prs)

    # 발표자 노트 부착 (content 슬라이드 인덱스)
    order = [None, None, None, "problem", None, "arch", None, "pipeline",
             "reliability", None, "ops", None, None]
    for idx, key in enumerate(order):
        if key:
            prs.slides[idx].notes_slide.notes_text_frame.text = NOTES[key]

    prs.save(out)
    print(f"saved: {out}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")


if __name__ == "__main__":
    main()
